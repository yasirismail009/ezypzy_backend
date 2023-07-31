from io import BytesIO
import os
import subprocess
from tempfile import NamedTemporaryFile
import traceback
from django.http import HttpResponse
from rest_framework.views import APIView
from ezypzy.serializers import FileTableSerializer
from utils.responses import created, internal_server_error, ok
from rest_framework.exceptions import ValidationError
import docx
from docx2pdf import convert
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import boto3
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate
from reportlab.platypus.tables import Table
from core import secrets
from django.core.files.base import ContentFile


class FileSave(APIView):
     def post(self, request):
        try:
            input_file = request.FILES.get('fileUpload')
            file_extension = os.path.splitext(input_file.name)[-1].lower()
            file_name = input_file.name
            if file_extension in [ ".pptx", ".ppt"]:
                    pdf_data = pptTopdf(input_file)
                    pdf_file = ContentFile(pdf_data, name=f"{file_name}.pdf")
                    data={
                    "fileType": file_extension,
                    "fileName": file_name,
                    "file": pdf_file,
                    "deviceId": "767871387198",
                    "converted": True
                }  
            elif file_extension in [ ".docx", ".doc"]:
                         pdf_data = docTopdf(input_file)
                         pdf_file = ContentFile(pdf_data, name=f"{file_name}.pdf")

                         data={
                           "fileType": file_extension,
                            "fileName": file_name,
                            "file": pdf_file,
                            "deviceId": "767871387198",
                            "converted": True
                        }  
            else:
                   data={
                           "fileType": file_extension,
                            "fileName": file_name,
                            "file": input_file,
                            "deviceId": "767871387198",
                            "converted": False
                        }
                 
            serilizeData=FileTableSerializer(data=data)
            if serilizeData.is_valid(raise_exception=True):
                    serilizeData.save()
                    return created(data=serilizeData.data)
            else:
                    return internal_server_error(message="error")

          
        except ValidationError as err:
            error_message = err.get_full_details()
            print(traceback.format_exc())
            return internal_server_error(message=error_message)


def pptTopdf(inputFile):
    presentation = Presentation(inputFile)

    # Initialize the PDF document
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    elements = []

    # Process each slide and extract content
    for slide in presentation.slides:
        content = []

        for shape in slide.shapes:
            if shape.has_text_frame:  # Corrected this line (removed parentheses)
                content.append(shape.text)

        # Add the content of each slide as a row in a table
        table = Table([content])
        table.setStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ])
        elements.append(table)

    # Build the PDF document
    doc.build(elements)

    # Get the PDF content from the buffer
    pdf_content = buffer.getvalue()

    # Close the buffer
    buffer.close()
    temp_docx_file = "temp_file.pdf"
    with open(temp_docx_file, 'wb') as temp_file:
        temp_file.write(pdf_content)

    # Return the ContentFile object as the response
    return pdf_content



def docTopdf(inputFile):
        file_content = inputFile.read()

        # Save the DOCX file to a temporary file on the server's filesystem
        temp_docx_file = "temp_file.docx"
        with open(temp_docx_file, 'wb') as temp_file:
            temp_file.write(file_content)

        # Convert DOCX content to PDF using docx2pdf
        convert(temp_docx_file)
        pdf_output_path = temp_docx_file.replace(".docx", ".pdf")

        with open(pdf_output_path, 'rb') as pdf_file:
            pdf_content = pdf_file.read()

        os.remove(temp_docx_file)
        os.remove(pdf_output_path)
        return pdf_content


  # Save the contents of the uploaded file to a temporary file
            # with NamedTemporaryFile(delete=False) as temp_file:
            #     for chunk in input_file.chunks():
            #         temp_file.write(chunk)

            # # Use the temporary file's path to open the presentation with aspose.slides
            # with slides.Presentation(temp_file.name) as presentation:
            #     presentation.save("presentation.pdf", slides.export.SaveFormat.PDF)

            # # Delete the temporary file
            # os.remove(temp_file.name)

            # # Return the PDF content as a response
            # response = HttpResponse("pdf_content", content_type='application/pdf')
            # response['Content-Disposition'] = 'attachment; filename="output.pdf"'
            # return response



#  class FileSave(APIView):
#     def post(self, request):
#         try:
#             input_file = request.FILES.get('fileUpload')
#             # Read the file content as bytes
#             presentation = Presentation(input_file)

#             # Initialize the PDF document
#             buffer = BytesIO()
#             doc = SimpleDocTemplate(buffer, pagesize=letter)
#             elements = []

#             # Process each slide and extract content
#             for slide in presentation.slides:
#                 content = []

#                 for shape in slide.shapes:
#                     if shape.has_text_frame:  # Corrected this line (removed parentheses)
#                         content.append(shape.text)

#                 # Add the content of each slide as a row in a table
#                 table = Table([content])
#                 table.setStyle([
#                     ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
#                     ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
#                     ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
#                     ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
#                     ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
#                     ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
#                     ('GRID', (0, 0), (-1, -1), 1, colors.black),
#                 ])
#                 elements.append(table)

#             # Build the PDF document
#             doc.build(elements)

#             # Get the PDF content from the buffer
#             pdf_content = buffer.getvalue()

#             # Close the buffer
#             buffer.close()
#             temp_docx_file = "temp_file.pdf"
#             with open(temp_docx_file, 'wb') as temp_file:
#                 temp_file.write(pdf_content)
#             # Return the PDF content as a response
#             response = HttpResponse(pdf_content, content_type='application/pdf')
#             response['Content-Disposition'] = 'attachment; filename="output.pdf"'
#             return response

#         except ValidationError as err:
#             error_message = err.get_full_details()
#             print(traceback.format_exc())
#             return internal_server_error(message=error_message)
       

# import os
# import traceback
# from django.http import HttpResponse
# from rest_framework.views import APIView
# from utils.responses import internal_server_error
# from rest_framework.exceptions import ValidationError
# from io import BytesIO
# import docx
# from docx2pdf import convert

# # Create your views here.

# class FileSave(APIView):
#     def post(self, request):
#         try:
#             input_file = request.FILES.get('fileUpload')
#             # Read the file content as bytes
#             file_content = input_file.read()

#             # Save the DOCX file to a temporary file on the server's filesystem
#             temp_docx_file = "temp_file.docx"
#             with open(temp_docx_file, 'wb') as temp_file:
#                 temp_file.write(file_content)

#             # Convert DOCX content to PDF using docx2pdf
#             pdf_output = convert(temp_docx_file)

#             # Clean up temporary file
#             os.remove(temp_docx_file)

#             # Serve the PDF to the user
#             response = HttpResponse(pdf_output, content_type='application/pdf')
#             response['Content-Disposition'] = 'attachment; filename="output.pdf"'

#             return response

#         except ValidationError as err:
#             error_message = err.get_full_details()
#             print(traceback.format_exc())
#             return internal_server_error(message=error_message)
